#!/usr/bin/env python3
"""
Build EA_Archimate_Examples.pptx: sample slides using local EMF stencils and
ArchiMate relationship arrows from this folder.

Run: python3 ea_archimate_emf/build_examples_presentation.py
Requires: python-pptx.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
EMF_ELEMENTS = ROOT / "emf"
EMF_REL = ROOT / "archimate_relations_emf" / "emf"
OUT = ROOT / "EA_Archimate_Examples.pptx"


def _pic(slide, rel_path: Path, left, top, width: Inches) -> None:
    if rel_path.is_file():
        slide.shapes.add_picture(str(rel_path), left, top, width=width)


def _banner(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.85))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True


def _caption(slide, text: str, left, top, width: Inches) -> None:
    box = slide.shapes.add_textbox(left, top, width, Inches(0.38))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.alignment = PP_ALIGN.CENTER


def _row_of_pictures(
    slide,
    base_dir: Path,
    entries: list[tuple[str, str]],
    *,
    left0: Inches,
    top: Inches,
    pic_w: Inches,
    gap: Inches,
    per_row: int,
) -> None:
    for idx, (fn, cap) in enumerate(entries):
        row, col = divmod(idx, per_row)
        left = left0 + col * (pic_w + gap)
        y = top + row * Inches(1.55)
        _pic(slide, base_dir / fn, left, y, pic_w)
        _caption(slide, cap, left, y + Inches(0.95), pic_w)


def main() -> int:
    if not EMF_ELEMENTS.is_dir():
        print(f"Missing element EMF folder: {EMF_ELEMENTS}", file=sys.stderr)
        return 1
    if not EMF_REL.is_dir():
        print(f"Missing relation EMF folder: {EMF_REL}", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    # --- Title ---
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Enterprise architecture in PowerPoint"
    s0.placeholders[1].text = (
        "Examples: ArchiMate-style elements, relationship connectors, and a small scenario.\n"
        "EMF files include four edge targets to align connectors."
    )

    # --- Business & application (grid) ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _banner(s1, "Layer examples — business & application elements")
    _row_of_pictures(
        s1,
        EMF_ELEMENTS,
        [
            ("Business_Actor.emf", "Business Actor"),
            ("Business_Role.emf", "Business Role"),
            ("Business_Process.emf", "Business Process"),
            ("Business_Service.emf", "Business Service"),
            ("Application_Component.emf", "Application Component"),
            ("Application_Service.emf", "Application Service"),
            ("Data_Object.emf", "Data Object"),
            ("Application_Function.emf", "Application Function"),
        ],
        left0=Inches(0.5),
        top=Inches(1.35),
        pic_w=Inches(2.05),
        gap=Inches(0.2),
        per_row=4,
    )

    # --- Technology, location, motivation ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _banner(s2, "Technology, location, and motivation")
    _row_of_pictures(
        s2,
        EMF_ELEMENTS,
        [
            ("Technology_Node.emf", "Node"),
            ("Technology_Interface.emf", "Technology Interface"),
            ("Technology.emf", "Technology"),
            ("Location.emf", "Location"),
            ("Representation.emf", "Representation"),
        ],
        left0=Inches(0.45),
        top=Inches(1.35),
        pic_w=Inches(2.15),
        gap=Inches(0.2),
        per_row=5,
    )
    _row_of_pictures(
        s2,
        EMF_ELEMENTS,
        [
            ("Constraint.emf", "Constraint"),
            ("Requirement.emf", "Requirement"),
            ("Goal.emf", "Goal"),
            ("Gap.emf", "Gap"),
        ],
        left0=Inches(0.45),
        top=Inches(3.15),
        pic_w=Inches(2.15),
        gap=Inches(0.2),
        per_row=4,
    )

    # --- Relationship connectors ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _banner(s3, "ArchiMate-style relationship connectors (EMF)")
    rels = [
        ("Archimate_Rel_Composition.emf", "Composition"),
        ("Archimate_Rel_Aggregation.emf", "Aggregation"),
        ("Archimate_Rel_Assignment.emf", "Assignment"),
        ("Archimate_Rel_Realization.emf", "Realization"),
        ("Archimate_Rel_Association.emf", "Association"),
        ("Archimate_Rel_Association_Directed.emf", "Assoc. (directed)"),
        ("Archimate_Rel_Serving.emf", "Serving"),
        ("Archimate_Rel_Access.emf", "Access"),
        ("Archimate_Rel_Influence.emf", "Influence"),
        ("Archimate_Rel_Triggering.emf", "Triggering"),
        ("Archimate_Rel_Flow.emf", "Flow"),
        ("Archimate_Rel_Specialization.emf", "Specialization"),
        ("Archimate_Rel_Junction.emf", "Junction"),
    ]
    _row_of_pictures(
        s3,
        EMF_REL,
        rels,
        left0=Inches(0.35),
        top=Inches(1.25),
        pic_w=Inches(1.85),
        gap=Inches(0.12),
        per_row=5,
    )

    # --- Mini scenario ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _banner(s4, "Mini scenario — actor, process, application service")
    y = Inches(1.85)
    w = Inches(2.0)
    rel_w = Inches(1.35)
    small = Inches(0.1)
    gap = Inches(0.22)
    x0 = Inches(0.55)
    _pic(s4, EMF_ELEMENTS / "Business_Actor.emf", x0, y, w)
    _caption(s4, "Actor", x0, y + Inches(1.0), w)

    x_rel_serving = x0 + w + small
    _pic(s4, EMF_REL / "Archimate_Rel_Serving.emf", x_rel_serving, y + Inches(0.35), rel_w)
    _caption(s4, "serving", x_rel_serving, y + Inches(1.05), rel_w)

    x_proc = x_rel_serving + rel_w + gap
    _pic(s4, EMF_ELEMENTS / "Business_Process.emf", x_proc, y, w)
    _caption(s4, "Business process", x_proc, y + Inches(1.0), w)

    x_rel_real = x_proc + w + small
    _pic(s4, EMF_REL / "Archimate_Rel_Realization.emf", x_rel_real, y + Inches(0.35), rel_w)
    _caption(s4, "realization", x_rel_real, y + Inches(1.05), rel_w)

    x_app = x_rel_real + rel_w + gap
    _pic(s4, EMF_ELEMENTS / "Application_Service.emf", x_app, y, w)
    _caption(s4, "Application service", x_app, y + Inches(1.0), w)

    note = s4.shapes.add_textbox(Inches(0.55), Inches(3.55), Inches(12.0), Inches(1.2))
    nf = note.text_frame
    nf.text = (
        "Insert native connectors from the Insert tab, or copy these horizontal relation EMFs "
        "and rotate/resize as needed. Blue rings on elements mark suggested left / right / top / bottom attachment guides."
    )
    nf.paragraphs[0].font.size = Pt(14)

    # --- Sources ---
    s5 = prs.slides.add_slide(prs.slide_layouts[1])
    s5.shapes.title.text = "Where the assets live"
    body = s5.shapes.placeholders[1].text_frame
    body.text = (
        "GitHub: fideocam / EAinPowerpoint\n\n"
        "• emf_icons — hand-built stencil set (with four edge anchor targets)\n"
        "• emf_exact_from_powerpoint — exports from Archimate_blank.pptx (same anchor treatment when regenerated)\n"
        "• emf_archimate_relations — canonical relationship lines\n\n"
        "Regenerate locally from the BlenderGPT workspace folder ea_archimate_emf/ "
        "(generate_icons.py, generate_archimate_relations.py, export_exact_from_ppt.py)."
    )
    for p in body.paragraphs:
        p.font.size = Pt(16)

    prs.save(OUT)
    print(f"Wrote {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
